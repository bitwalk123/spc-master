import pandas as pd
import numpy as np
import math
import re
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt


# _/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_
class ExcelSPC():
    filename = None
    sheets = None
    valid = False
    SL_flag = []

    header_master = ['Part Number', 'Description', 'Key Parameter', 'Parameter Name',
                     'LSL', 'Target', 'USL', 'Chart Type', 'Metrology', 'Multiple',
                     'Lower Tol', 'Upper Tol', 'Spec Type', 'CL Frozen',
                     'LCL', 'Avg', 'UCL', 'RLCL', 'R Avg', 'RUCL', 'CLCR Lower', 'CLCR Upper',
                     'Total # of Recent Points', '%OOC for Recent Points',
                     'Cpk for Recent Points', 'PPM for Recent Points',
                     'Parameter Classification', 'Product Classification',
                     'Recent Std Dev', 'Cpk for All Points', 'PPM for All Points',
                     'Cpk for Historic & Recent Points', 'PPM for Historic & Recent Points']

    def __init__(self, filename: str):
        self.filename: str = filename
        self.sheets: dict = self.read(filename)
        self.valid: bool = self.check_valid_sheet(self.sheets)
        if self.valid is False:
            return
        self.init_SL_flag()

    # -------------------------------------------------------------------------
    #  init_SL_flag
    #
    #  argument
    #    (none)
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def init_SL_flag(self):
        self.SL_flag = []
        df = self.get_master()
        n = len(df)
        for i in range(n):
            self.SL_flag.append(False)

    # -------------------------------------------------------------------------
    #  set_SL_flag
    #
    #  argument
    #    row  :
    #    flag :
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def set_SL_flag(self, row: int, flag: bool):
        self.SL_flag[row] = flag

    # -------------------------------------------------------------------------
    #  get_SL_flag
    #
    #  argument
    #    row
    #
    #  return
    #    bool: Spec Limit status for specified row
    # -------------------------------------------------------------------------
    def get_SL_flag(self, row: int) -> bool:
        return self.SL_flag[row]

    # -------------------------------------------------------------------------
    #  check_valid_sheet
    #  check if read file (sheets) has 'Master' tab
    #
    #  argument
    #    sheets : dataframe containing Excel contents
    #
    #  return
    #    True if dataframe is valid for SPC, otherwise False
    # -------------------------------------------------------------------------
    def check_valid_sheet(self, sheets) -> bool:
        # check if 'Master' tab exists
        if 'Master' in sheets.keys():
            if len(self.sheets['Master'].columns) == len(self.header_master):
                self.sheets['Master'].columns = self.header_master
                return True
            else:
                # TODO
                # need to identify extra empty column is added unintentionally
                print(self.sheets['Master'].columns)
                print(self.header_master)
                return False
        else:
            return False

    # -------------------------------------------------------------------------
    #  get_master
    #  get dataframe of 'Master' tab
    #
    #  argument
    #    (none)
    #
    #  return
    #    pandas dataframe of 'Master' tab
    # -------------------------------------------------------------------------
    def get_master(self) -> pd.DataFrame:
        df: pd.DataFrame = self.sheets['Master']
        # drop row if column 'Part Number' is NaN
        df = df.dropna(subset=['Part Number'])

        return df

    # -------------------------------------------------------------------------
    #  get_metrics
    #
    #  argument
    #    name_part :
    #    param     :
    #
    #  return
    #    dict - metrics dictionary for specified PART and PARAMETER
    # -------------------------------------------------------------------------
    def get_metrics(self, name_part, param):
        df = self.get_master()
        df1 = df[(df['Part Number'] == name_part) & (df['Parameter Name'] == param)]

        dict = {}
        dict['Part Number'] = list(df1['Part Number'])[0]
        dict['Description'] = list(df1['Description'])[0]
        dict['Key Parameter'] = list(df1['Key Parameter'])[0]
        dict['Parameter Name'] = list(df1['Parameter Name'])[0]
        dict['LSL'] = list(df1['LSL'])[0]
        dict['Target'] = list(df1['Target'])[0]
        dict['USL'] = list(df1['USL'])[0]
        dict['Chart Type'] = list(df1['Chart Type'])[0]
        dict['Metrology'] = list(df1['Metrology'])[0]
        dict['Multiple'] = list(df1['Multiple'])[0]
        dict['Lower Tol'] = list(df1['Lower Tol'])[0]
        dict['Upper Tol'] = list(df1['Upper Tol'])[0]
        dict['Spec Type'] = list(df1['Spec Type'])[0]
        dict['CL Frozen'] = list(df1['CL Frozen'])[0]
        dict['LCL'] = list(df1['LCL'])[0]
        dict['Avg'] = list(df1['Avg'])[0]
        dict['UCL'] = list(df1['UCL'])[0]
        dict['RLCL'] = list(df1['RLCL'])[0]
        dict['R Avg'] = list(df1['R Avg'])[0]
        dict['RUCL'] = list(df1['RUCL'])[0]
        dict['CLCR Lower'] = list(df1['CLCR Lower'])[0]
        dict['CLCR Upper'] = list(df1['CLCR Upper'])[0]
        dict['Total # of Recent Points'] = list(df1['Total # of Recent Points'])[0]
        dict['%OOC for Recent Points'] = list(df1['%OOC for Recent Points'])[0]
        dict['Cpk for Recent Points'] = list(df1['Cpk for Recent Points'])[0]
        dict['PPM for Recent Points'] = list(df1['PPM for Recent Points'])[0]
        dict['Recent Std Dev'] = list(df1['Recent Std Dev'])[0]
        dict['Cpk for All Points'] = list(df1['Cpk for All Points'])[0]
        dict['PPM for All Points'] = list(df1['PPM for All Points'])[0]
        dict['Cpk for Historic & Recent Points'] = list(df1['Cpk for Historic & Recent Points'])[0]
        dict['PPM for Historic & Recent Points'] = list(df1['PPM for Historic & Recent Points'])[0]

        return dict

    # -------------------------------------------------------------------------
    #  get_param_list
    #  get list of 'Parameter Name' of specified 'Part Number'
    #
    #  argument
    #    name_part : part name
    #
    #  return
    #    list of 'Parameter Name' of specified 'Part Number'
    # -------------------------------------------------------------------------
    def get_param_list(self, name_part):
        df = self.get_master()

        return list(df[df['Part Number'] == name_part]['Parameter Name'])

    # -------------------------------------------------------------------------
    #  get_part
    #  get dataframe of specified name_part tab
    #
    #  argument
    #    (none)
    #
    #  return
    #    pandas dataframe of specified name_part tab
    # -------------------------------------------------------------------------
    def get_part(self, name_part):
        # dataframe of specified name_part tab
        df = self.sheets[name_part]

        # delete row including NaN
        df = df.dropna(how='all')

        # _/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_
        #  the first row od data sheet is used for 'Create Charts' button for
        #  the Excel macro
        #
        #  So, new dataframe is created for this application
        # _/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_

        # obtain number of rows on this dataframe
        row_size = len(df)

        # extract data rows
        df1 = df[1:row_size]

        # extract column name used for this dataframe
        list_colname = list(df.loc[0])
        df1.columns = list_colname

        # for colname in ['Sample', 'Date', 'Job ID or Lot ID', 'Serial Number', 'Data Type']:
        df1 = df1.dropna(subset=['Data Type'])

        # eliminate 'Hide' data
        df2 = df1[df1['Data Type'] != 'Hide']

        return df2

    # -------------------------------------------------------------------------
    #  get_sheets
    #  get dataframe containing Excel contents
    #
    #  argument
    #    (none)
    #
    #  return
    #    array of pandas dataframe containing Excel tab/data
    # -------------------------------------------------------------------------
    def get_sheets(self):
        return self.sheets

    # -------------------------------------------------------------------------
    #  get_unique_part_list
    #  get unique part list found in 'Part Number' column in 'Master' tab
    #
    #  argument
    #    (none)
    #
    #  return
    #    list of unique 'Part Number'
    # -------------------------------------------------------------------------
    def get_unique_part_list(self):
        df = self.get_master()
        list_part = list(np.unique(df['Part Number']))

        return list_part

    # -------------------------------------------------------------------------
    #  get_header_master
    #  get header list used for making table
    #
    #  argument
    #    (none)
    #
    #  return
    #    self.header_master
    # -------------------------------------------------------------------------
    def get_header_master(self):
        return self.header_master

    # -------------------------------------------------------------------------
    #  read
    #  read specified Excel file
    #
    #  argument
    #    filename : Excel file
    #
    #  return
    #    array of pandas dataframe including all Excel sheets
    # -------------------------------------------------------------------------
    def read(self, filename):
        # read specified filename as Excel file including all tabs
        # return pd.read_excel(filename, sheet_name=None)

        df = pd.read_excel(
            filename,
            sheet_name=None,
            engine='openpyxl',
        )
        return df


# _/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_
class PowerPoint():
    ppt = None

    def __init__(self, template):
        # insert empty slide
        self.ppt = Presentation(template)

    # -------------------------------------------------------------------------
    #  add_slide
    #
    #  argument
    #    sheets :
    #    info   :
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def add_slide(self, sheets, info):
        metrics = sheets.get_metrics(info['PART'], info['PARAM'])
        metrics = self.check_dict(metrics)

        # ---------------------------------------------------------------------
        #  refer layout from original master
        # ---------------------------------------------------------------------
        slide_layout = self.ppt.slide_layouts[1]
        slide = self.ppt.slides.add_slide(slide_layout)
        shapes = slide.shapes

        # ---------------------------------------------------------------------
        #  slide title
        # ---------------------------------------------------------------------
        shapes.title.text = info['PART']

        # ---------------------------------------------------------------------
        # insert textbox
        # ---------------------------------------------------------------------
        ##### DEBUG ROUINE for PLACEHOLDER INDEX #####
        # for shape in slide.placeholders:
        #    print('%d %s' % (shape.placeholder_format.idx, shape.name))

        # Placeholder 1
        ph1 = shapes.placeholders[20]
        tf1 = ph1.text_frame
        tf1.text = self.get_body_text_1(metrics)

        # Placeholder 2
        ph2 = shapes.placeholders[21]
        tf2 = ph2.text_frame
        tf2.text = self.get_body_text_2(info, metrics)

        # ---------------------------------------------------------------------
        # insert image
        # ---------------------------------------------------------------------
        ileft = Inches(0)
        itop = Inches(1.92)
        iheight = Inches(3.5)

        slide.shapes.add_picture(info['IMAGE'], left=ileft, top=itop, height=iheight)

        # ---------------------------------------------------------------------
        # insert table
        # ---------------------------------------------------------------------
        # self.create_table(metrics, shapes)

    # -------------------------------------------------------------------------
    #  get_body_text_1
    #
    #  argument
    #    metrics :
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def get_body_text_1(self, metrics):
        # Chart Type Information
        if metrics['Chart Type'] == 'LJ':
            dist = 'Normal'
            ctype = 'Levey Jennings'
        elif metrics['Chart Type'] == 'IR':
            dist = 'Non-normal'
            ctype = 'IR'
        else:
            dist = 'Unknown'
            ctype = 'n/a'
        # Spec Information
        if metrics['Spec Type'] == 'Two-Sided':
            spec = '\tUSL = ' + metrics['USL'] + ', LSL = ' + metrics['LSL']
        elif metrics['Spec Type'] == 'One-Sided':
            spec = '\tUSL = ' + metrics['USL']
        else:
            spec = 'n/a'

        text = 'Inspection Method:\t' + metrics['Metrology'] \
               + '\tMeasurement Type:\t' + metrics['Multiple'] \
               + '\nDistribution:\t' + dist \
               + '\tParameter Type:\tKey' \
               + '\nChart Type:\t' + ctype \
               + spec
        return text

    # -------------------------------------------------------------------------
    #  get_body_text_2
    # -------------------------------------------------------------------------
    def get_body_text_2(self, info, metrics):
        if metrics['CL Frozen'] == 'Yes':
            SL_status = 'Frozen'
        else:
            SL_status = 'Not frozen'

        pattern = re.compile(r'(.*)\..*')  # left right side from floating point in mumber
        num = metrics['Total # of Recent Points']
        match = pattern.match(num)
        if match:
            num = match.group(1)

        text = 'Control Limit Status:\t' + SL_status \
               + '\nRecent Number of Data Points:\t' + num \
               + '\nDate of Last Lot Received:\t' + info['Date of Last Lot Received']

        return text

    # -------------------------------------------------------------------------
    #  create_table
    #
    #  argument
    #    metrics :
    #    shapes  :
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def create_table(self, metrics, shapes):
        rows = 10
        cols = 8
        top = Inches(4.4)
        left = Inches(0)
        width = Inches(10)
        height = Inches(0)
        table = shapes.add_table(rows, cols, left, top, width, height).table
        # set column widths
        table.columns[0].width = Inches(1.0)
        table.columns[1].width = Inches(3.5)
        table.columns[2].width = Inches(0.9)
        table.columns[3].width = Inches(0.7)
        table.columns[4].width = Inches(0.8)
        table.columns[5].width = Inches(0.7)
        table.columns[6].width = Inches(1.7)
        table.columns[7].width = Inches(0.7)
        table.cell(0, 0).text = 'index'
        table.cell(0, 1).text = 'value'
        table.cell(0, 2).text = 'index'
        table.cell(0, 3).text = 'value'
        table.cell(0, 4).text = 'index'
        table.cell(0, 5).text = 'value'
        table.cell(0, 6).text = 'index'
        table.cell(0, 7).text = 'value'
        table.cell(1, 0).text = 'Part Number'
        table.cell(2, 0).text = 'Description'
        table.cell(3, 0).text = 'Parameter Name'
        table.cell(4, 0).text = 'Key Parameter'
        table.cell(5, 0).text = 'Metrology'
        table.cell(1, 1).text = metrics['Part Number']
        table.cell(2, 1).text = metrics['Description']
        table.cell(3, 1).text = metrics['Parameter Name']
        table.cell(4, 1).text = metrics['Key Parameter']
        table.cell(5, 1).text = metrics['Metrology']
        table.cell(1, 2).text = 'LSL'
        table.cell(2, 2).text = 'Target'
        table.cell(3, 2).text = 'USL'
        table.cell(4, 2).text = 'Chart Type'
        table.cell(5, 2).text = 'Multiple'
        table.cell(6, 2).text = 'Spec Type'
        table.cell(1, 3).text = metrics['LSL']
        table.cell(2, 3).text = metrics['Target']
        table.cell(3, 3).text = metrics['USL']
        table.cell(4, 3).text = metrics['Chart Type']
        table.cell(5, 3).text = metrics['Multiple']
        table.cell(6, 3).text = metrics['Spec Type']
        table.cell(1, 4).text = 'CL Frozen'
        table.cell(2, 4).text = 'LCL'
        table.cell(3, 4).text = 'Avg'
        table.cell(4, 4).text = 'UCL'
        table.cell(5, 4).text = 'RLCL'
        table.cell(6, 4).text = 'R Avg'
        table.cell(7, 4).text = 'RUCL'
        table.cell(8, 4).text = 'CLCR Lower'
        table.cell(9, 4).text = 'CLCR Upper'
        table.cell(1, 5).text = metrics['CL Frozen']
        table.cell(2, 5).text = metrics['LCL']
        table.cell(3, 5).text = metrics['Avg']
        table.cell(4, 5).text = metrics['UCL']
        table.cell(5, 5).text = metrics['RLCL']
        table.cell(6, 5).text = metrics['R Avg']
        table.cell(7, 5).text = metrics['RUCL']
        table.cell(8, 5).text = metrics['CLCR Lower']
        table.cell(9, 5).text = metrics['CLCR Upper']
        table.cell(1, 6).text = 'Total # of Recent Points'
        table.cell(2, 6).text = '%OOC for Recent Points'
        table.cell(3, 6).text = 'Cpk for Recent Points'
        table.cell(4, 6).text = 'PPM for Recent Points'
        table.cell(5, 6).text = 'Recent Std Dev'
        table.cell(6, 6).text = 'Cpk for All Points'
        table.cell(7, 6).text = 'PPM for All Points'
        table.cell(8, 6).text = 'Cpk for Historic & Recent Points'
        table.cell(9, 6).text = 'PPM for Historic & Recent Points'
        table.cell(1, 7).text = metrics['Total # of Recent Points']
        table.cell(2, 7).text = metrics['%OOC for Recent Points']
        table.cell(3, 7).text = metrics['Cpk for Recent Points']
        table.cell(4, 7).text = metrics['PPM for Recent Points']
        table.cell(5, 7).text = metrics['Recent Std Dev']
        table.cell(6, 7).text = metrics['Cpk for All Points']
        table.cell(7, 7).text = metrics['PPM for All Points']
        table.cell(8, 7).text = metrics['Cpk for Historic & Recent Points']
        table.cell(9, 7).text = metrics['PPM for Historic & Recent Points']
        for r in range(rows):
            for c in range(cols):
                font = table.cell(r, c).text_frame.paragraphs[0].font
                font.size = Pt(7)

    # -------------------------------------------------------------------------
    #  check_dict
    #
    #  argument
    #    dict :
    #
    #  return
    #    dict after corrected
    # -------------------------------------------------------------------------
    def check_dict(self, dict):
        for key in dict:
            value = dict[key]
            if (type(value) is int):
                if math.isnan(value):
                    dict[key] = 'n/a'
                else:
                    dict[key] = str(value)

                continue

            if (type(value) is float):
                self.floatFormat(dict, key, value)

        return dict

    # -------------------------------------------------------------------------
    #  floatFormat
    #
    #  argument
    #    dict  :
    #    key   :
    #    value :
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def floatFormat(self, dict, key, value):
        if math.isnan(value):
            dict[key] = 'n/a'
        else:
            if abs(value) < 10:
                dict[key] = '{:.6f}'.format(value)
            else:
                n = int(math.log10(abs(value)))
                if n < 6:
                    f = '{:.' + str(6 - n) + 'f}'
                    dict[key] = f.format(value)
                else:
                    dict[key] = str(int(value))

    # -------------------------------------------------------------------------
    #  save PowerPoint file
    # -------------------------------------------------------------------------
    def save(self, save_path):
        self.ppt.save(save_path)

# ---
# PROGRAM END
