import numpy as np
import pathlib
import subprocess

from gi.repository import Gtk
from matplotlib.backends.backend_gtk3agg import FigureCanvasGTK3Agg as FigureCanvas
from matplotlib.figure import Figure

from pptx import Presentation
from pptx.util import Cm, Inches, Pt
from pptx.dml.color import RGBColor

from module import utils


# =============================================================================
#  CharWin - single chart window
# =============================================================================
class ChartWin(Gtk.Window):

    def __init__(self, info_master, widget, sheet):
        Gtk.Window.__init__(self)
        self.connect("delete-event", self.on_delete)
        self.info_master = info_master

        row = utils.register(int(widget.get_label()), 1, self.info_master.get_rows() - 1)
        self.info_master.select_row(row.get())

        # get Parameter Name & PART Number
        name_param, name_part = self.get_part_param(row, sheet)

        # _/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_
        #  HeaderBar
        self.hbar = Gtk.HeaderBar()
        self.hbar.set_show_close_button(True)
        self.set_hbar_title(name_part, name_param)

        # _/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_
        # Buttons at LEFT
        box_left = Gtk.Box(orientation=Gtk.Orientation.HORIZONTAL)
        self.hbar.pack_start(box_left)

        image_left = Gtk.Image.new_from_icon_name('go-previous', Gtk.IconSize.BUTTON)
        but_left = Gtk.Button()
        but_left.set_image(image_left)
        but_left.set_tooltip_text("go previous parameter")
        box_left.add(but_left)

        image_right = Gtk.Image.new_from_icon_name('go-next', Gtk.IconSize.BUTTON)
        but_right = Gtk.Button()
        but_right.set_image(image_right)
        but_right.set_tooltip_text("go next parameter")
        box_left.add(but_right)

        # _/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_
        # Buttons at RIGHT
        box_right = Gtk.Box(orientation=Gtk.Orientation.HORIZONTAL)
        self.hbar.pack_end(box_right)

        # configulation plot for chart
        image_ppt = Gtk.Image.new_from_icon_name('x-office-presentation', Gtk.IconSize.BUTTON)
        but_ppt = Gtk.Button()
        but_ppt.set_image(image_ppt)
        but_ppt.set_tooltip_text("create PowerPoint file")
        box_right.add(but_ppt)

        # configulation plot for chart
        image_conf = Gtk.Image.new_from_icon_name('applications-system', Gtk.IconSize.BUTTON)
        but_conf = Gtk.Button()
        but_conf.set_image(image_conf)
        but_conf.set_tooltip_text("chart setting")
        box_right.add(but_conf)

        # _/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_
        #  SPC chart
        # canvas = self.generate_spc_plot(sheet, name_part, name_param)
        self.box = Gtk.Box(orientation=Gtk.Orientation.VERTICAL)
        self.add(self.box)
        self.gen_chart(sheet, name_part, name_param)

        # _/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_/_
        #  binding for clicking on arrow button
        but_left.connect('clicked', self.on_arrow_left_clicked, row, sheet)
        but_right.connect('clicked', self.on_arrow_right_clicked, row, sheet)
        but_ppt.connect('clicked', self.on_ppt_clicked, sheet, name_part, name_param)

        self.show_all()

    # -------------------------------------------------------------------------
    #  set_hbar_title - set HeaderBar title to 'Part No' & 'Parameter Name'
    #
    #  argument
    #    name_part  : PART Number
    #    name_param : Parameter Name
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def set_hbar_title(self, name_part, name_param):
        self.hbar.props.title = name_part + " - " + name_param
        self.set_titlebar(self.hbar)

    # -------------------------------------------------------------------------
    #  gen_chart - generate chart
    #
    #  argument
    #    name_part  : PART Number
    #    name_param : Parameter Name
    #    sheet      : data sheet from Excel file
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def gen_chart(self, sheet, name_part, name_param):
        figure = make_trend_chart(sheet, name_part, name_param)
        canvas = FigureCanvas(figure)
        canvas.set_size_request(1500, 500)
        self.box.pack_start(canvas, expand=True, fill=True, padding=0)

    # -------------------------------------------------------------------------
    #  get_part_param - get PART No & Parameter Name from sheet
    #
    #  argument
    #    row   : row object on the Master Table
    #    sheet : data sheet from Excel file
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def get_part_param(self, row, sheet):
        df_master = sheet.get_master()
        df_row = (df_master.iloc[row.get() - 1])
        name_part = df_row['Part Number']
        name_param = df_row['Parameter Name']

        return name_param, name_part

    # -------------------------------------------------------------------------
    #  on_arrow_left_clicked - event handling when left arrow botton on
    #                          the HeaderBar is clicked.
    #
    #  argument
    #    widget : clicked widget
    #    row    : row object on the Master Table
    #    sheet  : data sheet from Excel file
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def on_arrow_left_clicked(self, widget, row, sheet):
        if row.isMin():
            return

        row.dec()
        self.update_chart(row, sheet)

    # -------------------------------------------------------------------------
    #  on_arrow_right_clicked - event handling when lright arrow botton on
    #                           the HeaderBar is clicked.
    #
    #  argument
    #    widget : clicked widget
    #    row    : row object on the Master Table
    #    sheet  : data sheet from Excel file
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def on_arrow_right_clicked(self, widget, row, sheet):
        if row.isMax():
            return

        row.inc()
        self.update_chart(row, sheet)

    # -------------------------------------------------------------------------
    #  update_chart - update chart
    #
    #  argument
    #    row    : row object on the Master Table
    #    sheet  : data sheet from Excel file
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def update_chart(self, row, sheet):
        # remove existence child widgets
        for child in self.box.get_children():
            self.box.remove(child)
            child.destroy()
        # add new chart
        self.info_master.select_row(row.get())
        name_param, name_part = self.get_part_param(row, sheet)
        self.set_hbar_title(name_part, name_param)
        self.gen_chart(name_part, name_param, sheet)
        self.show_all()

    def on_ppt_clicked(self, button, sheet, name_part, name_param):
        image_path = "./chart.png"
        template_path = "./template.pptx"
        save_path = "./output.pptx"

        # create chart
        figure = make_trend_chart(sheet, name_part, name_param)

        # create PNG file of plot
        figure.savefig(image_path)

        # create PowerPoint file
        gen_ppt(image_path, template_path, save_path)

        # open created file
        open_file_with_app(save_path)

        # complete messages
        #dialog = Gtk.MessageDialog(parent=self,
        #                           flags=0,
        #                           message_type=Gtk.MessageType.INFO,
        #                           buttons=Gtk.ButtonsType.OK,
        #                           text="generated PowerPoint file in " + save_path + ".")
        #dialog.run()
        #dialog.destroy()

    # -------------------------------------------------------------------------
    #  on_delete - event handling when close botton X on the window is clicked.
    #
    #  argument
    #    widget : clicked widget
    #    foo    : dummy
    #
    #  return
    #    (none)
    # -------------------------------------------------------------------------
    def on_delete(self, widget, foo):
        self.__del__()

    # -------------------------------------------------------------------------
    #  destructor
    #
    #  Note:
    #  Usually, destructor is not reauired in Python programing but this is
    #  prepared on purpose since special handling is required in this case.
    # -------------------------------------------------------------------------
    def __del__(self):
        # print('DEBUG')
        self.info_master.deselect_row()
        self.close()
        self.destroy()


# =============================================================================
#  GLOBAL FUNCTIONS
# =============================================================================
# -----------------------------------------------------------------------------
#  make_trend_chart
# -----------------------------------------------------------------------------
def make_trend_chart(sheet, name_part, name_param):
    metrics = sheet.get_metrics(name_part, name_param)
    df = sheet.get_part(name_part)
    x = df['Sample']
    y = df[name_param]
    fig = Figure(dpi=100, figsize=(10, 3.5))
    splot = fig.add_subplot(111, title=name_param, ylabel='Value')
    splot.grid(True)

    if metrics['Spec Type'] == 'Two-Sided':
        if not np.isnan(metrics['USL']):
            splot.axhline(y=metrics['USL'], linewidth=1, color='blue', label='USL')
        if not np.isnan(metrics['UCL']):
            splot.axhline(y=metrics['UCL'], linewidth=1, color='red', label='UCL')
        if not np.isnan(metrics['Target']):
            splot.axhline(y=metrics['Target'], linewidth=1, color='purple', label='Target')
        if not np.isnan(metrics['LCL']):
            splot.axhline(y=metrics['LCL'], linewidth=1, color='red', label='LCL')
        if not np.isnan(metrics['LSL']):
            splot.axhline(y=metrics['LSL'], linewidth=1, color='blue', label='LSL')
    elif metrics['Spec Type'] == 'One-Sided':
        if not np.isnan(metrics['USL']):
            splot.axhline(y=metrics['USL'], linewidth=1, color='blue', label='USL')
        if not np.isnan(metrics['UCL']):
            splot.axhline(y=metrics['UCL'], linewidth=1, color='red', label='UCL')

    # Avg
    splot.axhline(y=metrics['Avg'], linewidth=1, color='green', label='Avg')
    # Line
    splot.plot(x, y, linewidth=1, color="gray")

    size_oos = 60
    size_ooc = 100
    if metrics['Spec Type'] == 'Two-Sided':
        # OOC check
        x_ooc = x[(df[name_param] < metrics['LCL']) | (df[name_param] > metrics['UCL'])]
        y_ooc = y[(df[name_param] < metrics['LCL']) | (df[name_param] > metrics['UCL'])]
        splot.scatter(x_ooc, y_ooc, s=size_ooc, c='orange', marker='o', label="Recent")
        # OOS check
        x_oos = x[(df[name_param] < metrics['LSL']) | (df[name_param] > metrics['USL'])]
        y_oos = y[(df[name_param] < metrics['LSL']) | (df[name_param] > metrics['USL'])]
        splot.scatter(x_oos, y_oos, s=size_oos, c='red', marker='o', label="Recent")
    elif metrics['Spec Type'] == 'One-Sided':
        # OOC check
        x_ooc = x[(df[name_param] > metrics['UCL'])]
        y_ooc = y[(df[name_param] > metrics['UCL'])]
        splot.scatter(x_ooc, y_ooc, s=size_ooc, c='orange', marker='o', label="Recent")
        # OOS check
        x_oos = x[(df[name_param] > metrics['USL'])]
        y_oos = y[(df[name_param] > metrics['USL'])]
        splot.scatter(x_oos, y_oos, s=size_oos, c='red', marker='o', label="Recent")
    splot.scatter(x, y, s=20, c='black', marker='o', label="Recent")
    x_label = splot.get_xlim()[1]
    if metrics['Spec Type'] == 'Two-Sided':
        if not np.isnan(metrics['USL']):
            splot.text(x_label, y=metrics['USL'], s=' USL', color='blue')
        if not np.isnan(metrics['UCL']):
            splot.text(x_label, y=metrics['UCL'], s=' UCL', color='red')
        if not np.isnan(metrics['Target']):
            splot.text(x_label, y=metrics['Target'], s=' Target', color='purple')
        if not np.isnan(metrics['LCL']):
            splot.text(x_label, y=metrics['LCL'], s=' LCL', color='red')
        if not np.isnan(metrics['LSL']):
            splot.text(x_label, y=metrics['LSL'], s=' LSL', color='blue')
    elif metrics['Spec Type'] == 'One-Sided':
        if not np.isnan(metrics['USL']):
            splot.text(x_label, y=metrics['USL'], s=' USL', color='blue')
        if not np.isnan(metrics['UCL']):
            splot.text(x_label, y=metrics['UCL'], s=' UCL', color='red')

    # Avg
    splot.text(x_label, y=metrics['Avg'], s=' Avg', color='green')

    return fig


def gen_ppt(image_path, template_path, save_path):
    # insert empty slide
    presentation = Presentation(template_path)
    # refer layout from original master
    title_slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    # slide title
    slide_title = "SPC chart (Example)"
    shapes.title.text = slide_title

    # ---------------------------------------------------
    # insert image
    # ---------------------------------------------------
    # insert position
    pic_left = Inches(0)
    pic_top = Inches(0.84)
    # image height
    pic_height = Inches(3.5)
    slide.shapes.add_picture(image_path, pic_left, pic_top, height=pic_height)

    # ---------------------------------------------------------------------
    #  save PowerPoint file
    # ---------------------------------------------------------------------
    presentation.save(save_path)

# -------------------------------------------------------------------------
#  open_file_with_app
#
#  argument
#    name_file :  file to open
# -------------------------------------------------------------------------
def open_file_with_app(name_file):
    link_file = pathlib.PurePath(name_file)
    # Explorer can cover all cases on Windows NT
    subprocess.Popen(['explorer', link_file])


# ---
# PROGRAM END
